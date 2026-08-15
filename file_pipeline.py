# =========================================================
# [file_pipeline.py] 
# 문서, 이미지, 미디어 파일 데이터 추출/전처리 및 
# 로컬 AI(Ollama) 메타데이터 생성 모듈
# =========================================================

# 1. 파이썬 표준 라이브러리 (기능별)
import os          # 파일 경로 확인, 파일 크기 검사, 확장자 추출 모듈
import re          # 정규표현식(특수문자 제거, 패턴 매칭) 모듈
import json        # JSON 문자열 변환 및 데이터 파싱 모듈
import io          # 메모리 내 바이너리 바이트 버퍼 처리 모듈
import zlib        # HWP 파일 데이터 압축 해제(Decompress) 모듈
import zipfile     # HWPX/DOCX 등 ZIP 포맷 압축 해제 모듈
import xml.etree.ElementTree as ET  # XML 구조 파일 텍스트 추출용 모듈
import requests    # 로컬 AI(Ollama) HTTP API 통신용 라이브러리
from datetime import datetime       # 타임스탬프(분석 시간) 기록용 모듈
from typing import Dict, Any, Tuple # 파이썬 함수 리턴 타입 명시용 모듈

# 2. 문서 및 이미지 파싱용 외부 제3자 라이브러리
from pypdf import PdfReader                  # PDF 파일 텍스트 추출
from docx import Document                    # Word(.docx) 파일 문단 텍스트 추출
import openpyxl                              # Excel(.xlsx) 셀 데이터 추출
from PIL import Image, UnidentifiedImageError # 이미지 리사이징 및 손상 검사

# 3. 외부 선택 설치 패키지 안전검사 (try-except 모듈 동적 로딩)
# olefile (구버전 HWP 바이너리 파싱용)
try:
    import olefile
    HAS_OLEFILE = True
except ImportError:
    HAS_OLEFILE = False

# python-pptx (PowerPoint 파싱용)
try:
    from pptx import Presentation
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False

# openai-whisper (음성 인식 STT 전용)
try:
    import whisper
    HAS_WHISPER = True
except ImportError:
    HAS_WHISPER = False


# =========================================================
# [커스텀 예외 클래스]
# =========================================================
class FilePreprocessError(Exception):
    """파일 전처리(읽기/해석) 과정에서 오류가 발생했을 때 던지는 예외 클래스"""
    pass


# =========================================================
# [Step 1] 파일 원문/데이터 추출 클래스 (TextExtractor)
# =========================================================
class TextExtractor:
    """확장자별 문서, 이미지, 오디오/비디오 데이터 추출 및 6중 방어막 예외 처리 클래스"""

    # A. 처리 가능한 파일 확장자 그룹 정의
    IMAGE_EXTENSIONS = ('.jpg', '.jpeg', '.png', '.webp', '.bmp', '.gif', '.tiff', '.tif')
    
    DOC_EXTENSIONS = (
        '.txt', '.pdf', '.docx', '.xlsx', '.pptx', 
        '.hwp', '.hwpx',
        '.csv', '.json', '.xml', '.yaml', '.yml', 
        '.html', '.htm', '.md', '.markdown'
    )
    
    AUDIO_VIDEO_EXTENSIONS = ('.mp3', '.wav', '.m4a', '.mp4', '.mkv', '.avi')

    # B. 보안 및 오류 방지를 위해 명시적으로 시스템 입력을 차단할 압축 파일
    ARCHIVE_EXTENSIONS = ('.zip', '.rar', '.7z', '.tar', '.gz', '.bz2', '.iso')

    def __init__(
        self, 
        max_chars: int = 2000, 
        max_img_size: int = 512, 
        whisper_model_name: str = "base",
        max_doc_size_mb: int = 50,       # 문서 파일 최대 50MB 용량 제한
        max_img_size_mb: int = 100,      # 이미지 파일 최대 100MB 용량 제한
        max_media_size_mb: int = 500     # 오디오/비디오 최대 500MB 용량 제한
    ):
        """기본 제한 설정(글자 수 제한, 최대 용량 제한 등)을 초기화하는 생성자"""
        self.max_chars = max_chars
        self.max_img_size = max_img_size
        self.whisper_model_name = whisper_model_name
        
        # MB 바이트 단위로 계산 변환 (1MB = 1024 * 1024 Bytes)
        self.max_doc_bytes = max_doc_size_mb * 1024 * 1024
        self.max_img_bytes = max_img_size_mb * 1024 * 1024
        self.max_media_bytes = max_media_size_mb * 1024 * 1024
        
        self._whisper_model = None  # Whisper 모델 지연 로딩(필요할 때 메모리 적재)

    def _sanitize_text(self, text: str) -> str:
        """문서 내부의 깨진 문자나 시스템 특수 제어 문자를 안전하게 제거하는 정제 함수"""
        if not text:
            return ""
        return re.sub(r'[\x00-\x08\x0b\x0c\x0e-\x1f]', '', text)

    def is_image_file(self, file_path: str) -> bool:
        """해당 파일 경로가 지원되는 이미지 확장자인지 여부 확인"""
        ext = os.path.splitext(file_path)[1].lower()
        return ext in self.IMAGE_EXTENSIONS

    def is_media_file(self, file_path: str) -> bool:
        """해당 파일 경로가 지원되는 오디오/비디오 확장자인지 여부 확인"""
        ext = os.path.splitext(file_path)[1].lower()
        return ext in self.AUDIO_VIDEO_EXTENSIONS

    # ---------------------------------------------------------
    # [1] 이미지 파일 전처리 함수 (애니메이션 GIF 및 확장자 불일치 예외 완벽 방어)
    # ---------------------------------------------------------
    def process_image(self, file_path: str) -> Tuple[bytes, str]:
        """이미지 파일 유효성 검사, 썸네일 축소 및 메모리 바이너리(Bytes) 변환 함수"""
        
        # 1-1. 존재 여부 검사
        if not os.path.exists(file_path):
            return b"", "ERROR: 존재하지 않는 이미지 파일입니다."

        # 1-2. 0 바이트 빈 파일 검사
        file_size = os.path.getsize(file_path)
        if file_size == 0:
            return b"", "ERROR: 내용이 없는 빈 이미지 파일(0 Byte)입니다."

        # 1-3. 용량 제한 초과 검사
        if file_size > self.max_img_bytes:
            max_mb = self.max_img_bytes // (1024 * 1024)
            return b"", f"ERROR: 이미지 용량이 제한({max_mb}MB)을 초과했습니다. ({file_size / (1024*1024):.1f}MB)"

        try:
            # 💡 [버그 수정] img.verify() 이중 호출 제거: verify()를 부르면 파일 포인터가 손상되어
            # 특정 WEBP, GIF, 확장자가 변형된 이미지 파일에서 UnidentifiedImageError가 터지던 문제를 완벽 해결했습니다.
            with Image.open(file_path) as img:
                # 애니메이션 GIF 등의 다중 프레임 이미지일 경우 첫 번째 대표 프레임 선택
                if getattr(img, "is_animated", False):
                    img.seek(0)

                img_rgb = img.convert("RGB") # Vision AI 전달용 RGB 포맷 통합
                img_rgb.thumbnail((self.max_img_size, self.max_img_size)) # 축소 리사이징

                buffer = io.BytesIO()
                img_rgb.save(buffer, format="JPEG", quality=80) # JPEG 압축 저장
                return buffer.getvalue(), "SUCCESS"

        except PermissionError:
            return b"", "ERROR: 이미지 파일 접근 권한이 없거나 다른 프로그램에서 사용 중입니다."
        except Image.DecompressionBombError:
            return b"", "ERROR: 해상도가 너무 큰 비정상적인 이미지(DecompressionBomb)입니다."
        except Exception as e:
            # 가짜 확장자나 헤더 손상 시에도 전체 스캔이 중단되지 않도록 예외 차단 후 메시지 반환
            return b"", f"ERROR: 이미지 읽기 및 변환 실패 ({str(e)})"

    # ---------------------------------------------------------
    # [2] 오디오/비디오 미디어 파일 전처리 (FFmpeg 미설치/STT 오류 안전 방어)
    # ---------------------------------------------------------
    def process_media(self, file_path: str) -> Tuple[str, str]:
        """Whisper AI 모델을 활용해 음성을 텍스트로 변환하는 함수"""
        file_name = os.path.basename(file_path)
        
        # 2-1. 패키지 설치 여부 검사
        if not HAS_WHISPER:
            return f"미디어 파일: {file_name} (openai-whisper 미설치)", "SUCCESS"

        # 2-2. 파일 존재 여부 및 용량 검사
        if not os.path.exists(file_path) or os.path.getsize(file_path) == 0:
            return f"미디어 파일: {file_name} (빈 파일 또는 존재하지 않음)", "SUCCESS"

        try:
            # Whisper 모델 지연 로딩
            if self._whisper_model is None:
                self._whisper_model = whisper.load_model(self.whisper_model_name)

            # STT 음성 추출 실행
            result = self._whisper_model.transcribe(file_path)
            extracted_text = result.get("text", "").strip()

            if not extracted_text:
                extracted_text = f"미디어 파일명: {file_name} (음성 인식 데이터 없음)"

            clean_text = self._sanitize_text(extracted_text)
            return clean_text[:self.max_chars].strip(), "SUCCESS"

        except Exception as e:
            # FFmpeg가 미설치되어 있거나 디코딩이 실패해도 파일명을 텍스트로 넘겨 AI 분석 진행!
            return f"음성 추출 불가 미디어 파일: {file_name} (STT 실패: {str(e)})", "SUCCESS"

    # ---------------------------------------------------------
    # [3] 일반 문서 및 데이터 파일 텍스트 추출 메인 함수 (파싱 오류 완벽 방어)
    # ---------------------------------------------------------
    def extract(self, file_path: str) -> Tuple[str, str]:
        """확장자별 개별 추출 알고리즘을 호출하고 예외 상황을 일괄 제어하는 메인 함수"""
        
        if not os.path.exists(file_path):
            return "", "ERROR: 존재하지 않는 파일입니다."

        file_size = os.path.getsize(file_path)
        if file_size == 0:
            return "", "ERROR: 내용이 없는 빈 파일(0 Byte)입니다."

        ext = os.path.splitext(file_path)[1].lower()

        # 방어막 1: 단순 압축 파일 명시적 거부
        if ext in self.ARCHIVE_EXTENSIONS:
            return "", f"ERROR: 압축 파일({ext})은 지원하지 않습니다. 압축을 해제한 후 개별 파일로 업로드해 주세요."

        # 방어막 2: 문서 용량 제한 검사
        if file_size > self.max_doc_bytes:
            max_mb = self.max_doc_bytes // (1024 * 1024)
            return "", f"ERROR: 문서 용량이 제한({max_mb}MB)을 초과했습니다. ({file_size / (1024*1024):.1f}MB)"

        # 방어막 3: 구버전 MS 오피스 파일(.doc, .xls, .ppt) 지원 제외
        if ext in ['.doc', '.xls', '.ppt']:
            return "", f"ERROR: 구버전 오피스 파일({ext})은 지원하지 않습니다. 최신 포맷(.docx, .xlsx, .pptx)으로 변환해 주세요."

        try:
            # 포맷별 개별 파서 분기 연결
            if ext in ['.txt', '.csv', '.json', '.xml', '.yaml', '.yml', '.html', '.htm', '.md', '.markdown']:
                text = self._read_txt(file_path)
            elif ext == '.pdf':
                text = self._read_pdf(file_path)
            elif ext == '.docx':
                text = self._read_docx(file_path)
            elif ext == '.xlsx':
                text = self._read_xlsx(file_path)
            elif ext == '.pptx':
                text = self._read_pptx(file_path)
            elif ext == '.hwpx':
                text = self._read_hwpx(file_path)
            elif ext == '.hwp':
                text = self._read_hwp(file_path)
            else:
                return "", f"ERROR: 지원하지 않는 파일 확장자입니다 ({ext})"

            clean_text = self._sanitize_text(text)

            # 💡 텍스트 추출이 빈 값이어도 파일명 자체를 전달하여 AI 태깅 수행
            if not clean_text or not clean_text.strip():
                clean_text = f"문서 파일명: {os.path.basename(file_path)} (내부 텍스트 내용 없음)"

            return clean_text[:self.max_chars].strip(), "SUCCESS"

        except Exception as e:
            # 💡 [버그 수정] HWP/PDF/DOCX 등 파싱 중 FilePreprocessError 등의 예외가 터져도
            # 백엔드 스레드가 멈추지 않고 파일명으로 대체하여 AI 태깅이 수월하게 진행되도록 최상단에서 방어합니다.
            file_name = os.path.basename(file_path)
            return f"문서 파일명: {file_name} (내부 데이터 해석 실패: {str(e)})", "SUCCESS"

    # --- 포맷별 하위 파싱 메서드 모음 ---

    def _read_txt(self, path: str) -> str:
        """대용량 TXT 방어: 최대 1MB 스트리밍 읽기 (UTF-8 / CP949 자동 적용)"""
        try:
            with open(path, 'r', encoding='utf-8') as f:
                return f.read(1024 * 1024)
        except UnicodeDecodeError:
            try:
                with open(path, 'r', encoding='cp949', errors='ignore') as f:
                    return f.read(1024 * 1024)
            except Exception as e:
                raise FilePreprocessError(f"텍스트 인코딩 읽기 실패: {str(e)}")

    def _read_pdf(self, path: str) -> str:
        """PDF 문서의 각 페이지에서 텍스트 읽기 및 암호화 여부 체크"""
        try:
            reader = PdfReader(path)
            if reader.is_encrypted:
                raise FilePreprocessError("암호로 보호된 PDF 파일입니다.")

            extracted_text = []
            total_len = 0

            for page in reader.pages:
                page_text = page.extract_text() or ""
                if page_text.strip():
                    extracted_text.append(page_text)
                    total_len += len(page_text)

                if total_len >= self.max_chars:
                    break

            return "\n".join(extracted_text)
        except FilePreprocessError:
            raise
        except Exception as e:
            raise FilePreprocessError(f"PDF 파싱 실패 ({str(e)})")

    def _read_docx(self, path: str) -> str:
        """Word 문서(.docx)의 각 문단 Paragraph 텍스트 추출"""
        try:
            doc = Document(path)
            full_text = [p.text for p in doc.paragraphs if p.text.strip()]
            return "\n".join(full_text)
        except Exception as e:
            raise FilePreprocessError(f"Word 파일 파싱 실패 ({str(e)})")

    def _read_xlsx(self, path: str) -> str:
        """Excel(.xlsx) 메모리 절약 모드(read_only)로 시트 내 셀 데이터 읽기"""
        wb = None
        try:
            wb = openpyxl.load_workbook(path, data_only=True, read_only=True)
            text_list = []
            current_len = 0

            for sheet in wb.worksheets:
                for row in sheet.iter_rows(values_only=True):
                    row_str = " ".join([str(cell) for cell in row if cell is not None])
                    if row_str.strip():
                        text_list.append(row_str)
                        current_len += len(row_str)

                    if current_len >= self.max_chars:
                        break
                if current_len >= self.max_chars:
                    break

            return "\n".join(text_list)
        except Exception as e:
            raise FilePreprocessError(f"Excel 파일 파싱 실패 ({str(e)})")
        finally:
            if wb:
                wb.close()

    def _read_pptx(self, path: str) -> str:
        """PowerPoint(.pptx) 슬라이드 도형 내 텍스트 파싱"""
        if not HAS_PPTX:
            raise FilePreprocessError("python-pptx 패키지가 설치되지 않았습니다. ('pip install python-pptx' 필요)")

        try:
            prs = Presentation(path)
            text_list = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text_list.append(shape.text.strip())
            return "\n".join(text_list)
        except Exception as e:
            raise FilePreprocessError(f"PPTX 파싱 실패 ({str(e)})")

    def _read_hwpx(self, path: str) -> str:
        """HWPX(ZIP 압축 해제) 형태의 Contents/section.xml 본문 추출"""
        try:
            text_list = []
            with zipfile.ZipFile(path, 'r') as zf:
                section_files = sorted([f for f in zf.namelist() if re.match(r'Contents/section\d+\.xml', f)])
                
                for sec_file in section_files:
                    xml_data = zf.read(sec_file)
                    root = ET.fromstring(xml_data)
                    
                    for elem in root.iter():
                        if elem.tag.endswith('t') and elem.text:
                            text_list.append(elem.text.strip())
                            
            return "\n".join(text_list)
        except zipfile.BadZipFile:
            raise FilePreprocessError("올바른 HWPX 압축 파일이 아니거나 손상되었습니다.")
        except Exception as e:
            raise FilePreprocessError(f"HWPX 파싱 실패 ({str(e)})")

    def _read_hwp(self, path: str) -> str:
        """HWP 오피스 바이너리 문서 OLE 구조에서 BodyText 비트 스트림 복호화 및 추출"""
        if not HAS_OLEFILE:
            raise FilePreprocessError("olefile 패키지가 설치되지 않았습니다. ('pip install olefile' 필요)")

        try:
            ole = olefile.OleFileIO(path)
            dirs = ole.listdir()

            body_sections = [d for d in dirs if d[0] == 'BodyText']
            body_sections.sort()

            if not body_sections:
                ole.close()
                raise FilePreprocessError("HWP 문서 내 본문(BodyText) 영역을 찾을 수 없습니다.")

            text_list = []
            for section in body_sections:
                stream = ole.openstream(section).read()
                
                # Deflate 압축 데이터 해제
                try:
                    decompressed = zlib.decompress(stream, -15)
                except zlib.error:
                    decompressed = stream

                # Record 구조 단위 오프셋 탐색
                i = 0
                while i < len(decompressed):
                    if i + 4 > len(decompressed):
                        break
                    record_header = int.from_bytes(decompressed[i:i+4], 'little')
                    tag_id = record_header & 0x3FF
                    size = (record_header >> 20) & 0xFFF
                    
                    if size == 0xFFF:
                        if i + 8 > len(decompressed):
                            break
                        size = int.from_bytes(decompressed[i+4:i+8], 'little')
                        i += 8
                    else:
                        i += 4

                    if i + size > len(decompressed):
                        break

                    payload = decompressed[i:i+size]
                    i += size

                    if tag_id == 67: # HWPTAG_PARA_TEXT 문단 태그 ID
                        para_text = payload.decode('utf-16le', errors='ignore')
                        clean_para = re.sub(r'[\x00-\x09\x0b-\x1f]', '', para_text).strip()
                        if clean_para:
                            text_list.append(clean_para)

            ole.close()

            if not text_list:
                raise FilePreprocessError("HWP 문서에서 추출 가능한 텍스트가 없습니다.")

            return "\n".join(text_list)

        except FilePreprocessError:
            raise
        except Exception as e:
            raise FilePreprocessError(f"HWP 파싱 실패 (암호화되었거나 손상된 파일일 수 있습니다: {str(e)})")


# =========================================================
# [Step 2] 로컬 AI(Ollama) 메타데이터 생성 클래스 (FileAnalyzer)
# =========================================================
class FileAnalyzer:
    """추출된 원문을 기반으로 로컬 LLM/Vision 모델에 요청하여 메타데이터 JSON을 만드는 클래스"""

    def __init__(
        self, 
        ollama_url: str = "http://localhost:11434", 
        text_model: str = "qwen2.5:3b",
        vision_model: str = "llava"
    ):
        """Ollama API URL 및 사용할 텍스트/비전 LLM 모델명 초기화"""
        self.ollama_api_url = f"{ollama_url.rstrip('/')}/api/generate"
        self.ollama_chat_url = f"{ollama_url.rstrip('/')}/api/chat"
        self.text_model = text_model
        self.vision_model = vision_model

    def _get_file_info(self, file_path: str) -> Dict[str, Any]:
        """파일의 원본 이름, 확장자, 바이트 크기 및 분석 시각 메타데이터 구성"""
        file_name = os.path.basename(file_path)
        ext = os.path.splitext(file_name)[1].lower()
        size = os.path.getsize(file_path) if os.path.exists(file_path) else 0
        
        return {
            "original_name": file_name,
            "file_extension": ext,
            "file_size_bytes": size,
            "analyzed_at": datetime.now().isoformat()
        }

    # ---------------------------------------------------------
    # 텍스트 기반 문서 분석 및 메타데이터 JSON 생성
    # ---------------------------------------------------------
    def analyze_document_text(self, file_path: str, extracted_text: str) -> Dict[str, Any]:
        """문서 텍스트 원문을 로컬 텍스트 LLM에 전달하여 제목, 태그, 요약 JSON을 추출하는 함수"""
        file_name = os.path.basename(file_path)
        file_info = self._get_file_info(file_path)

        prompt = f"""
You are a professional file metadata analyzer. Analyze the provided text content and generate structured metadata in JSON format.

[File Information]
- Original Filename: {file_name}
- Content Text:
{extracted_text}

[Output Requirements]
Return ONLY a valid JSON object with the following keys:
1. "display_name": A clean, concise, and descriptive title for the file in Korean (Do NOT include file extension).
2. "tags": An array of 3 to 5 relevant keyword strings (without '#' symbol).
3. "description": A brief 1-2 sentence summary of the content in Korean.

Example JSON output format:
{{
  "display_name": "JSL 일본어 초급 교재 1권",
  "tags": ["JSL", "일본어", "초급", "교재"],
  "description": "JSL 일본어 초급 학습용 문법 및 단어 교재입니다."
}}
"""

        payload = {
            "model": self.text_model,
            "prompt": prompt,
            "format": "json", # Ollama에 JSON 출력 형식 강제 설정
            "stream": False,
            "options": {
                "temperature": 0.2,
                "num_predict": 400
            }
        }

        try:
            # Ollama 서버 API 호출
            response = requests.post(self.ollama_api_url, json=payload, timeout=90)
            response.raise_for_status()

            res_data = response.json()
            raw_response_text = res_data.get("response", "").strip()

            # 정규표현식으로 순수 JSON 영역만 파싱
            match = re.search(r'\{.*\}', raw_response_text, re.DOTALL)
            json_str = match.group(0) if match else raw_response_text
            parsed_json = json.loads(json_str)

            tags = parsed_json.get("tags", [])
            desc = parsed_json.get("description", "")
            tags_formatted = ", ".join([f"#{t}" for t in tags]) if tags else "#일반"
            ai_comment_str = f"태그: {tags_formatted} / 코멘트: {desc}"

            # ✨ [@TYPE: @DB 추가] MainProcessor 순서도 라우팅과 규격 통일
            return {
                "@TYPE": "@DB",
                "status": "SUCCESS",
                "file_info": file_info,
                "metadata": {
                    "@TYPE": "@DB",
                    "display_name": parsed_json.get("display_name", file_info["original_name"].rsplit('.', 1)[0]),
                    "tags": tags,
                    "description": desc,
                    "ai_comment": ai_comment_str, # GUI 디스플레이용 완성형 문자열
                    "ocr_text": ""
                },
                "error": None
            }

        except requests.exceptions.ConnectionError:
            return self._build_fallback_response(file_info, "Ollama AI 서버에 연결할 수 없습니다. (Ollama 실행 필요)")
        except requests.exceptions.Timeout:
            return self._build_fallback_response(file_info, "AI 분석 시간 초과 (Timeout - 응답 지연)")
        except json.JSONDecodeError:
            return self._build_fallback_response(file_info, "AI 응답 파싱 실패 (유효하지 않은 JSON 구조)")
        except Exception as e:
            return self._build_fallback_response(file_info, f"Text AI 분석 중 예외 발생 ({str(e)})")

    # ---------------------------------------------------------
    # 이미지 바이너리 기반 Vision 멀티모달 분석
    # ---------------------------------------------------------
    def analyze_image_bytes(self, file_path: str, img_bytes: bytes) -> Dict[str, Any]:
        """이미지 바이트 데이터를 Base64로 인코딩하여 Vision 모델(llava)에 전달하는 함수"""
        import base64
        
        file_name = os.path.basename(file_path)
        file_info = self._get_file_info(file_path)

        # 이미지를 Base64 문자열로 변환
        base64_img = base64.b64encode(img_bytes).decode('utf-8')

        prompt = f"""
You are an image analysis expert. Analyze the provided image and return a JSON object with:
1. "display_name": A descriptive name for this image in Korean (without extension).
2. "tags": An array of 3 to 5 relevant keyword strings in Korean (without '#' symbol).
3. "description": A brief 1-2 sentence description/summary of what is shown in the image in Korean.

Filename: {file_name}

Example JSON output format:
{{
  "display_name": "팀 회의 화이트보드 메모",
  "tags": ["회의", "아이디어", "아키텍처", "일정"],
  "description": "팀 프로젝트 아키텍처 및 세부 일정이 기록된 화이트보드 이미지입니다."
}}
"""

        payload = {
            "model": self.vision_model,
            "messages": [
                {
                    "role": "user",
                    "content": prompt,
                    "images": [base64_img]
                }
            ],
            "format": "json",
            "stream": False,
            "options": {
                "temperature": 0.2,
                "num_predict": 200
            },
            "keep_alive": "5m"
        }

        try:
            response = requests.post(self.ollama_chat_url, json=payload, timeout=180)
            response.raise_for_status()

            res_data = response.json()
            raw_response_text = res_data.get("message", {}).get("content", "").strip()

            match = re.search(r'\{.*\}', raw_response_text, re.DOTALL)
            json_str = match.group(0) if match else raw_response_text
            parsed_json = json.loads(json_str)

            tags = parsed_json.get("tags", [])
            desc = parsed_json.get("description", "")
            tags_formatted = ", ".join([f"#{t}" for t in tags]) if tags else "#이미지"
            ai_comment_str = f"태그: {tags_formatted} / 코멘트: {desc}"

            # ✨ [@TYPE: @DB 추가] MainProcessor 순서도 라우팅과 규격 통일
            return {
                "@TYPE": "@DB",
                "status": "SUCCESS",
                "file_info": file_info,
                "metadata": {
                    "@TYPE": "@DB",
                    "display_name": parsed_json.get("display_name", file_info["original_name"].rsplit('.', 1)[0]),
                    "tags": tags,
                    "description": desc,
                    "ai_comment": ai_comment_str, # GUI 디스플레이용 완성형 문자열
                    "ocr_text": ""
                },
                "error": None
            }

        except requests.exceptions.ConnectionError:
            return self._build_fallback_response(file_info, "Ollama AI 서버에 연결할 수 없습니다. (Ollama 실행 필요)")
        except requests.exceptions.Timeout:
            return self._build_fallback_response(file_info, "Vision AI 분석 시간 초과 (Timeout - 로딩 또는 연산 지연)")
        except json.JSONDecodeError:
            return self._build_fallback_response(file_info, "Vision AI 응답 파싱 실패 (유효하지 않은 JSON 구조)")
        except Exception as e:
            return self._build_fallback_response(file_info, f"Vision AI 분석 중 예외 발생 ({str(e)})")

    # ---------------------------------------------------------
    # 예외 발생 시 안전하게 기본값을 채워주는 폴백(Fallback) 함수
    # ---------------------------------------------------------
    def _build_fallback_response(self, file_info: Dict[str, Any], error_message: str) -> Dict[str, Any]:
        """AI 분석 실패 시 프로그램 멈춤 없이 최소 메타데이터로 구성된 실패 응답 JSON 반환"""
        default_name = file_info["original_name"].rsplit('.', 1)[0]
        return {
            "@TYPE": "@DB",
            "status": "FAILED",
            "file_info": file_info,
            "metadata": {
                "@TYPE": "@DB",
                "display_name": default_name,
                "tags": [],
                "description": f"분석 실패: {error_message}",
                "ai_comment": f"#분석실패 / 코멘트: {error_message}",
                "ocr_text": ""
            },
            "error": error_message
        }


# =====================================================================
# 통합 단독 테스트 실행부 (main)
# =====================================================================
if __name__ == "__main__":
    # 파서 및 분석기 객체 인스턴스 생성
    extractor = TextExtractor(max_chars=2000, max_img_size=512)
    analyzer = FileAnalyzer(text_model="qwen2.5:3b", vision_model="llava")

    # 테스트할 파일 샘플 리스트
    test_files = [
        "sample.txt",
        "sample.pdf",
        "sample.docx",
        "sample.xlsx",
        "sample.pptx",
        "sample.hwp",
        "sample.hwpx",
        "sample.json",
        "sample.png",
        "sample.gif",
        "sample.mp3",
        "sample.mp4",
        "sample.zip"  # 압축 파일 거부 차단 동작 테스트
    ]

    print("=== [통합 테스트] 문서/이미지/미디어 파일 전처리 및 AI 분석 ===")

    for file_path in test_files:
        print(f"\n========================================")
        print(f"📄 대상 파일: {file_path}")

        if not os.path.exists(file_path):
            print(f"[경고] '{file_path}' 파일을 찾을 수 없습니다. 테스트를 스킵합니다.")
            continue

        # 🖼️ 1. 이미지 파일인 경우의 분기 처리
        if extractor.is_image_file(file_path):
            img_bytes, status = extractor.process_image(file_path)
            print(f"[Step 1 이미지 전처리 상태]: {status}")

            if status == "SUCCESS":
                print(f"[Step 2 Vision AI({analyzer.vision_model}) 분석 요청 중...]")
                ai_result = analyzer.analyze_image_bytes(file_path, img_bytes)
                print("[Step 2 AI 분석 결과 (JSON)]:")
                print(json.dumps(ai_result, ensure_ascii=False, indent=2))
            else:
                print(f"[Step 1 예외/거부 사유]: {status}")

        # 🎵 2. 오디오/비디오 미디어 파일인 경우의 분기 처리
        elif extractor.is_media_file(file_path):
            extracted_text, status = extractor.process_media(file_path)
            print(f"[Step 1 음성 인식(STT) 상태]: {status}")

            if status == "SUCCESS":
                print(f"[Step 1 추출 텍스트 (최대 300자)]:\n---")
                print(extracted_text[:300] + ("..." if len(extracted_text) > 300 else ""))
                print("---")

                print(f"[Step 2 Text AI({analyzer.text_model}) 분석 요청 중...]")
                ai_result = analyzer.analyze_document_text(file_path, extracted_text)
                print("[Step 2 AI 분석 결과 (JSON)]:")
                print(json.dumps(ai_result, ensure_ascii=False, indent=2))
            else:
                print(f"[Step 1 예외/거부 사유]: {status}")

        # 📝 3. 일반 문서 파일인 경우의 분기 처리
        else:
            extracted_text, status = extractor.extract(file_path)
            print(f"[Step 1 텍스트 추출 상태]: {status}")

            if status == "SUCCESS":
                print(f"[Step 1 추출 텍스트 (최대 300자)]:\n---")
                print(extracted_text[:300] + ("..." if len(extracted_text) > 300 else ""))
                print("---")

                print(f"[Step 2 Text AI({analyzer.text_model}) 분석 요청 중...]")
                ai_result = analyzer.analyze_document_text(file_path, extracted_text)
                print("[Step 2 AI 분석 결과 (JSON)]:")
                print(json.dumps(ai_result, ensure_ascii=False, indent=2))
            else:
                print(f"[Step 1 예외/거부 사유]: {status}")